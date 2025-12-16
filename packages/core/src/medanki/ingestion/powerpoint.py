"""PowerPoint extraction for MedAnki."""

from __future__ import annotations

from pathlib import Path
from typing import Protocol, runtime_checkable

from medanki.ingestion.errors import IngestionError
from medanki.models.document import Document, Section
from medanki.models.enums import ContentType


@runtime_checkable
class IPowerPointExtractor(Protocol):
    """Protocol for PowerPoint extraction."""

    async def extract(self, path: Path) -> Document:
        """Extract content from a PowerPoint file."""
        ...


class PowerPointExtractor:
    """Extractor for PowerPoint (.pptx) files.

    Extracts text from slides and speaker notes, preserving slide order.
    Note: Legacy .ppt format is not supported.
    """

    async def extract(self, path: Path) -> Document:
        """Extract content from a PowerPoint file.

        Args:
            path: Path to the PowerPoint file.

        Returns:
            Document with extracted text and slide-based sections.

        Raises:
            IngestionError: If file is missing, corrupted, or legacy .ppt format.
        """
        if not path.exists():
            raise IngestionError(
                f"File not found: {path}",
                path=str(path),
            )

        suffix = path.suffix.lower()
        if suffix == ".ppt":
            raise IngestionError(
                f"Legacy .ppt format is not supported. Please convert to .pptx: {path}",
                path=str(path),
            )

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise IngestionError(
                "python-pptx is required for PowerPoint extraction. "
                "Install with: pip install python-pptx",
                path=str(path),
            ) from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise IngestionError(
                f"Failed to open PowerPoint file: {path}. File may be corrupted.",
                path=str(path),
            ) from exc

        sections: list[Section] = []
        all_text_parts: list[str] = []
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            slide_title = f"Slide {slide_num}"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_texts.append(text)

                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        try:
                            from pptx.enum.shapes import PP_PLACEHOLDER

                            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                                slide_title = text
                        except (AttributeError, ImportError):
                            if slide_texts and not slide_title.startswith("Slide"):
                                pass
                            elif slide_texts:
                                slide_title = slide_texts[0]

                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text and para_text not in slide_texts:
                            slide_texts.append(para_text)

            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_texts.append(f"[Notes: {notes_text}]")
            except Exception:
                pass

            if slide_texts:
                if slide_title == f"Slide {slide_num}" and slide_texts:
                    slide_title = slide_texts[0][:50]
                    if len(slide_texts[0]) > 50:
                        slide_title += "..."

                section_content = "\n".join(slide_texts)
                sections.append(
                    Section(
                        title=slide_title,
                        content=section_content,
                        level=1,
                        page_numbers=[slide_num],
                    )
                )
                all_text_parts.extend(slide_texts)

        raw_text = "\n\n".join(all_text_parts)

        return Document(
            source_path=path,
            content_type=ContentType.POWERPOINT_SLIDES,
            raw_text=raw_text,
            sections=sections,
            metadata={"slide_count": slide_count},
        )
