"""Tests for PowerPoint extraction."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import pytest

if TYPE_CHECKING:
    from medanki.ingestion.powerpoint import PowerPointExtractor


@pytest.fixture
def powerpoint_extractor() -> PowerPointExtractor:
    """Create a PowerPointExtractor instance."""
    from medanki.ingestion.powerpoint import PowerPointExtractor

    return PowerPointExtractor()


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a sample .pptx file for testing."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    slide_layout = prs.slide_layouts[5]
    slide1 = prs.slides.add_slide(slide_layout)
    title_shape = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Introduction to Cell Biology"

    content_shape = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_shape.text_frame
    content_frame.text = "Cells are the basic unit of life."
    p = content_frame.add_paragraph()
    p.text = "The cell membrane controls what enters and exits."

    slide2 = prs.slides.add_slide(slide_layout)
    title_shape2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_shape2.text_frame.text = "Mitochondria"

    content_shape2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_shape2.text_frame.text = "The mitochondria is the powerhouse of the cell."

    slide3 = prs.slides.add_slide(slide_layout)
    title_shape3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_shape3.text_frame.text = "Summary"

    pptx_path = tmp_path / "test_presentation.pptx"
    prs.save(str(pptx_path))

    return pptx_path


@pytest.fixture
def empty_pptx(tmp_path: Path) -> Path:
    """Create an empty .pptx file for testing."""
    from pptx import Presentation

    prs = Presentation()
    pptx_path = tmp_path / "empty.pptx"
    prs.save(str(pptx_path))

    return pptx_path


@pytest.fixture
def pptx_with_notes(tmp_path: Path) -> Path:
    """Create a .pptx file with speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_shape.text_frame.text = "Slide with Notes"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are the speaker notes with additional details."

    pptx_path = tmp_path / "with_notes.pptx"
    prs.save(str(pptx_path))

    return pptx_path


class TestPowerPointExtraction:
    """Tests for PowerPoint extraction functionality."""

    @pytest.mark.asyncio
    async def test_extracts_text_from_slides(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Extracts text content from all slides."""
        document = await powerpoint_extractor.extract(sample_pptx)

        assert "Introduction to Cell Biology" in document.raw_text
        assert "Cells are the basic unit of life" in document.raw_text
        assert "mitochondria is the powerhouse" in document.raw_text
        assert "Summary" in document.raw_text

    @pytest.mark.asyncio
    async def test_preserves_slide_order(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Text is extracted in slide order."""
        document = await powerpoint_extractor.extract(sample_pptx)

        intro_pos = document.raw_text.find("Introduction to Cell Biology")
        mito_pos = document.raw_text.find("Mitochondria")
        summary_pos = document.raw_text.find("Summary")

        assert intro_pos < mito_pos < summary_pos

    @pytest.mark.asyncio
    async def test_creates_sections_per_slide(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Each slide becomes a section."""
        document = await powerpoint_extractor.extract(sample_pptx)

        assert len(document.sections) == 3
        assert document.sections[0].title == "Introduction to Cell Biology"
        assert document.sections[1].title == "Mitochondria"
        assert document.sections[2].title == "Summary"

    @pytest.mark.asyncio
    async def test_sections_have_slide_numbers(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Sections include page numbers (slide numbers)."""
        document = await powerpoint_extractor.extract(sample_pptx)

        assert document.sections[0].page_numbers == [1]
        assert document.sections[1].page_numbers == [2]
        assert document.sections[2].page_numbers == [3]

    @pytest.mark.asyncio
    async def test_handles_empty_presentation(
        self,
        powerpoint_extractor: PowerPointExtractor,
        empty_pptx: Path,
    ) -> None:
        """Empty presentation returns empty document."""
        document = await powerpoint_extractor.extract(empty_pptx)

        assert document.raw_text.strip() == ""
        assert len(document.sections) == 0

    @pytest.mark.asyncio
    async def test_extracts_speaker_notes(
        self,
        powerpoint_extractor: PowerPointExtractor,
        pptx_with_notes: Path,
    ) -> None:
        """Speaker notes are included in extraction."""
        document = await powerpoint_extractor.extract(pptx_with_notes)

        assert "speaker notes with additional details" in document.raw_text

    @pytest.mark.asyncio
    async def test_sets_correct_content_type(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Document has POWERPOINT_SLIDES content type."""
        from medanki.models.enums import ContentType

        document = await powerpoint_extractor.extract(sample_pptx)

        assert document.content_type == ContentType.POWERPOINT_SLIDES

    @pytest.mark.asyncio
    async def test_sets_source_path(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Document has correct source path."""
        document = await powerpoint_extractor.extract(sample_pptx)

        assert document.source_path == sample_pptx

    @pytest.mark.asyncio
    async def test_stores_metadata(
        self,
        powerpoint_extractor: PowerPointExtractor,
        sample_pptx: Path,
    ) -> None:
        """Metadata includes slide count."""
        document = await powerpoint_extractor.extract(sample_pptx)

        assert document.metadata.get("slide_count") == 3

    @pytest.mark.asyncio
    async def test_missing_file_raises_error(
        self,
        powerpoint_extractor: PowerPointExtractor,
        tmp_path: Path,
    ) -> None:
        """Missing file raises IngestionError."""
        from medanki.ingestion.errors import IngestionError

        missing_file = tmp_path / "nonexistent.pptx"

        with pytest.raises(IngestionError):
            await powerpoint_extractor.extract(missing_file)

    @pytest.mark.asyncio
    async def test_corrupted_file_raises_error(
        self,
        powerpoint_extractor: PowerPointExtractor,
        tmp_path: Path,
    ) -> None:
        """Corrupted file raises IngestionError."""
        from medanki.ingestion.errors import IngestionError

        corrupted_file = tmp_path / "corrupted.pptx"
        corrupted_file.write_bytes(b"not a real pptx file")

        with pytest.raises(IngestionError):
            await powerpoint_extractor.extract(corrupted_file)


class TestLegacyPptFormat:
    """Tests for legacy .ppt format handling."""

    @pytest.mark.asyncio
    async def test_ppt_format_raises_helpful_error(
        self,
        powerpoint_extractor: PowerPointExtractor,
        tmp_path: Path,
    ) -> None:
        """Legacy .ppt format raises error with helpful message."""
        from medanki.ingestion.errors import IngestionError

        legacy_file = tmp_path / "legacy.ppt"
        legacy_file.write_bytes(b"fake ppt content")

        with pytest.raises(IngestionError) as exc_info:
            await powerpoint_extractor.extract(legacy_file)

        assert ".pptx" in str(exc_info.value).lower() or "convert" in str(exc_info.value).lower()
